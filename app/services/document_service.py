"""
Document Processing Service
Handles PDF and PPTX file processing
"""
import logging
import json
import re
from pathlib import Path
from typing import List
import fitz  # PyMuPDF
import pptx
from openai import OpenAI

from app.core.config import settings

logger = logging.getLogger(__name__)


class DocumentService:
    """Service for processing documents and extracting questions"""
    
    def __init__(self):
        self.client = OpenAI(api_key=settings.OPENAI_API_KEY)
    
    def _extract_text_from_pdf(self, pdf_path: Path, max_chars: int = 2000) -> str:
        """Extract text from PDF file"""
        text = ""
        with fitz.open(pdf_path) as doc:
            for page in doc:
                text += page.get_text()
                if len(text) > max_chars:
                    break
        return text[:max_chars]
    
    def _extract_text_from_ppt(self, ppt_path: Path, max_chars: int = 2000) -> str:
        """Extract text from PowerPoint file"""
        text = ""
        prs = pptx.Presentation(ppt_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    if len(text) > max_chars:
                        break
        return text[:max_chars]
    
    def _extract_json_from_text(self, text: str) -> dict:
        """Extract JSON from text response"""
        try:
            match = re.search(r"\{.*\}", text, re.DOTALL)
            return json.loads(match.group()) if match else json.loads(text)
        except Exception:
            return None
    
    def _ask_openai(self, prompt: str) -> str:
        """Query OpenAI API"""
        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
        )
        return response.choices[0].message.content
    
    async def process_document(self, file_path: Path) -> List[str]:
        """
        Process document and extract questions
        """
        try:
            # Extract text based on file type
            if file_path.suffix.lower() == '.pdf':
                text = self._extract_text_from_pdf(file_path)
            elif file_path.suffix.lower() == '.pptx':
                text = self._extract_text_from_ppt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_path.suffix}")
            
            # Generate questions using OpenAI
            prompt = f"""
            Generate 5 Yes/No questions from this text.
            Return JSON only:
            {{ "questions": ["..."] }}

            Text:
            {text}
            """
            
            raw_response = self._ask_openai(prompt)
            data = self._extract_json_from_text(raw_response)
            
            if data and "questions" in data:
                return data["questions"]
            else:
                raise ValueError("Failed to extract questions from response")
            
        except Exception as e:
            logger.error(f"Document processing failed: {str(e)}", exc_info=True)
            raise
